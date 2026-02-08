"""Slide copying utilities."""

from copy import deepcopy
from pptx.slide import Slide
from pptx import Presentation
from pptx.opc.constants import RELATIONSHIP_TYPE as RT


class SlideCopier:
    """Handles copying slides between presentations."""

    @staticmethod
    def copy_slide(source_prs: Presentation, source_slide_index: int, target_prs: Presentation) -> Slide:
        """
        Copy a slide from source presentation to target presentation.

        This method creates a deep copy of the slide including all shapes,
        images, formatting, and layout information using deepcopy at the XML level.

        Args:
            source_prs: Source presentation
            source_slide_index: Index of slide to copy (0-based)
            target_prs: Target presentation

        Returns:
            The newly created slide in target presentation
        """
        source_slide = source_prs.slides[source_slide_index]

        # Ensure target presentation has the same slide size as source
        SlideCopier._copy_slide_size(source_prs, target_prs)

        # Get the source slide's layout relationship ID
        # This is the most reliable way to match layouts, especially for custom layouts
        source_slide_part = source_slide.part
        target_prs_part = target_prs.part

        # Find the layout by matching the slide master and layout relationships
        slide_layout = None
        source_layout = source_slide.slide_layout

        try:
            # Get the rId of the slide layout from the source slide
            source_layout_rId = None
            for rel_id, rel in source_slide_part.rels.items():
                if rel.target_part == source_layout.part:
                    source_layout_rId = rel_id
                    break

            # Get the layout part's relationship to its slide master
            source_layout_part = source_layout.part
            source_master = source_layout.slide_master

            # In target presentation, find the matching slide master
            target_master = None
            for master in target_prs.slide_masters:
                if master.name == source_master.name:
                    target_master = master
                    break

            if target_master is None and len(target_prs.slide_masters) > 0:
                # Use first slide master as fallback
                target_master = target_prs.slide_masters[0]

            # Find the matching layout in the target master
            if target_master:
                # Try to match by layout name within the same master
                for layout in target_master.slide_layouts:
                    if layout.name == source_layout.name:
                        slide_layout = layout
                        break

                # If no match by name, try by index within master's layouts
                if slide_layout is None:
                    source_master_layouts = list(source_master.slide_layouts)
                    target_master_layouts = list(target_master.slide_layouts)

                    for i, layout in enumerate(source_master_layouts):
                        if layout.name == source_layout.name:
                            if i < len(target_master_layouts):
                                slide_layout = target_master_layouts[i]
                            break

        except (AttributeError, IndexError) as e:
            pass

        # Final fallback: use any available layout
        if slide_layout is None:
            if len(target_prs.slide_layouts) > 6:
                slide_layout = target_prs.slide_layouts[6]
            else:
                slide_layout = target_prs.slide_layouts[0]

        # Create new slide with the layout
        dest_slide = target_prs.slides.add_slide(slide_layout)

        # Copy all shapes using deepcopy at XML level
        for shape in source_slide.shapes:
            try:
                new_element = deepcopy(shape.element)
                dest_slide.shapes._spTree.insert_element_before(new_element, "p:extLst")
            except Exception as e:
                # If deepcopy fails for a shape, try to continue with others
                continue

        # Copy image relationships (pictures)
        SlideCopier._copy_images(source_slide, dest_slide)

        return dest_slide

    @staticmethod
    def _copy_slide_size(source_prs: Presentation, target_prs: Presentation):
        """
        Copy slide dimensions from source to target presentation.

        Args:
            source_prs: Source presentation
            target_prs: Target presentation
        """
        try:
            target_prs.slide_width = source_prs.slide_width
            target_prs.slide_height = source_prs.slide_height
        except Exception:
            # If setting slide size fails, continue
            pass

    @staticmethod
    def _copy_images(source_slide: Slide, dest_slide: Slide):
        """
        Copy image parts and relationships from source slide to destination slide.

        This ensures that images referenced in the copied shapes are available
        in the target presentation with the correct relationship IDs.

        Args:
            source_slide: Source slide
            dest_slide: Destination slide
        """
        try:
            from io import BytesIO

            # Get the parts
            source_part = source_slide.part
            dest_part = dest_slide.part

            # Create a mapping of old rId to new rId for images
            rId_mapping = {}

            # Iterate through all relationships in the source slide
            for rel_id, rel in source_part.rels.items():
                # Check if this is an image relationship
                if rel.reltype == RT.IMAGE:
                    # Get the image part
                    image_part = rel.target_part
                    image_blob = image_part.blob

                    # Wrap bytes in BytesIO to make it file-like
                    image_stream = BytesIO(image_blob)

                    # Add the image to destination and get the new relationship
                    # In python-pptx 1.0.2+, this returns a tuple: (image_part, rId)
                    result = dest_part.get_or_add_image_part(image_stream)

                    # Handle both old and new API
                    if isinstance(result, tuple):
                        # New API (python-pptx 1.0.2+): returns (image_part, rId)
                        new_image_part, new_rId = result
                        rId_mapping[rel_id] = new_rId
                        continue
                    else:
                        # Old API: returns image_part only
                        new_image_part = result

                    # For old API, find the relationship ID for this image
                    for new_rel_id, new_rel in dest_part.rels.items():
                        if (new_rel.reltype == RT.IMAGE and
                            new_rel.target_part == new_image_part):
                            rId_mapping[rel_id] = new_rel_id
                            break

            # Update all image shape relationships in the destination slide
            # This is necessary because deepcopy copied the old rIds
            if rId_mapping:
                # Define namespaces
                p_ns = 'http://schemas.openxmlformats.org/presentationml/2006/main'
                a_ns = 'http://schemas.openxmlformats.org/drawingml/2006/main'
                r_ns = 'http://schemas.openxmlformats.org/officeDocument/2006/relationships'

                # Find all picture shapes using findall
                pics = dest_slide.element.findall(f'.//{{{p_ns}}}pic')
                for pic in pics:
                    # Look for the blip element which contains the image reference
                    blips = pic.findall(f'.//{{{a_ns}}}blip')
                    for blip in blips:
                        # Get the embed attribute (r:embed)
                        embed_attr = f'{{{r_ns}}}embed'
                        old_rId = blip.get(embed_attr)

                        if old_rId and old_rId in rId_mapping:
                            # Update to the new rId
                            blip.set(embed_attr, rId_mapping[old_rId])

        except Exception as e:
            # If copying images fails, continue (the slide is already copied)
            pass
