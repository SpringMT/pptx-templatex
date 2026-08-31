"""Main template engine for processing PowerPoint files."""

import copy
import io
import json
import re
from pathlib import Path
from typing import Any, Dict, Union

import lxml.etree as etree
from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.dml import MSO_COLOR_TYPE
from pptx.opc.constants import RELATIONSHIP_TYPE as RT
from pptx.opc.package import Part
from pptx.oxml.ns import qn
from pptx.slide import Slide
from pptx.util import Pt
from pptx_slide_copier import SlideCopier

from .exceptions import TemplateError
from .placeholder_replacer import PlaceholderReplacer
from .rich_text_parser import TextSegment, is_rich_text, parse_rich_text

# `{{ img:key }}` 画像マーカー。key は英数・. - _ [ ] を許可（nested/index キーに合わせる）。
_IMG_MARKER = re.compile(r"\{\{\s*img:([\w.\-\[\]]+)\s*\}\}")

# PowerPoint 2019 / Microsoft 365 がベクター描画に使う svgBlip 拡張
_SVG_NS = "http://schemas.microsoft.com/office/drawing/2016/SVG/main"
_SVG_EXT_URI = "{96DAC541-7B7A-43D3-8B79-37D633B846F1}"

# SVGフォールバックPNGのラスタライズ倍率（配置枠に対して）。
# 旧PowerPoint・Keynote・Google Slides等はフォールバック側で表示されるため高めにしておく。
_SVG_FALLBACK_SCALE = 3
_EMU_PER_PX = 914400 // 96  # 96dpi換算


class TemplateEngine:
    """
    PowerPoint template engine that copies slides and replaces placeholders.

    Usage:
        engine = TemplateEngine("template.pptx")
        config = {
            "slides": [
                {"src_page": 1, "replace_texts": {"name": "John", "age": "30"}},
                {"src_page": 2, "replace_texts": {"items": [{"name": "Item 1"}]}}
            ]
        }
        engine.process(config, "output.pptx")
    """

    def __init__(self, template_path: Union[str, Path]):
        """
        Initialize the template engine with a source PowerPoint file.

        Args:
            template_path: Path to the template PPTX file

        Raises:
            TemplateError: If template file cannot be loaded
        """
        self.template_path = Path(template_path)
        if not self.template_path.exists():
            raise TemplateError(f"Template file not found: {template_path}")

        # Store the path but don't load the presentation yet
        # We'll load it fresh for each process() call to avoid state issues
        try:
            # Just validate that the file can be opened
            test_prs = Presentation(str(self.template_path))
            self.template_prs = test_prs
        except Exception as e:
            raise TemplateError(f"Failed to load template: {str(e)}")

    def _copy_slide(self, source_slide_index: int, target_prs: Presentation) -> Slide:
        """
        Copy a slide from the template to the target presentation.

        Args:
            source_slide_index: Index of the slide to copy from the template (0-based)
            target_prs: The target presentation to copy to

        Returns:
            The newly created slide in the target presentation
        """
        return SlideCopier.copy_slide(self.template_prs, source_slide_index, target_prs)

    # ------------------------------------------------------------------
    # Rich-text helpers
    # ------------------------------------------------------------------

    @staticmethod
    def _save_run_fmt(run):
        """Return a dict of formatting attributes from a run."""
        fmt: Dict[str, Any] = {
            "name": run.font.name,
            "size": run.font.size,
            "bold": run.font.bold,
            "italic": run.font.italic,
            "underline": run.font.underline,
            "color_type": None,
            "color_rgb": None,
            "color_scheme_xml": None,
        }
        try:
            if hasattr(run.font.color, "type"):
                fmt["color_type"] = run.font.color.type
                if fmt["color_type"] == MSO_COLOR_TYPE.RGB:
                    fmt["color_rgb"] = run.font.color.rgb
                elif fmt["color_type"] is not None:
                    # schemeClr など RGB以外のカラー指定: XMLノードをそのまま保存
                    rPr = run._r.find(qn("a:rPr"))
                    if rPr is not None:
                        solidFill = rPr.find(qn("a:solidFill"))
                        if solidFill is not None:
                            fmt["color_scheme_xml"] = copy.deepcopy(solidFill)
        except Exception:
            pass
        return fmt

    @staticmethod
    def _apply_run_fmt(run, fmt: Dict[str, Any]):
        """Apply saved formatting dict to a run."""
        if fmt["name"] is not None:
            run.font.name = fmt["name"]
        if fmt["size"] is not None:
            run.font.size = fmt["size"]
        if fmt["bold"] is not None:
            run.font.bold = fmt["bold"]
        if fmt["italic"] is not None:
            run.font.italic = fmt["italic"]
        if fmt["underline"] is not None:
            run.font.underline = fmt["underline"]
        if fmt["color_type"] == MSO_COLOR_TYPE.RGB and fmt["color_rgb"] is not None:
            try:
                run.font.color.rgb = fmt["color_rgb"]
            except Exception:
                pass
        elif fmt["color_scheme_xml"] is not None:
            # schemeClr などをXMLレベルで復元
            try:
                rPr = run._r.find(qn("a:rPr"))
                if rPr is None:
                    rPr = etree.SubElement(run._r, qn("a:rPr"))
                    run._r.insert(0, rPr)
                existing = rPr.find(qn("a:solidFill"))
                if existing is not None:
                    rPr.remove(existing)
                rPr.insert(0, copy.deepcopy(fmt["color_scheme_xml"]))
            except Exception:
                pass

    @staticmethod
    def _apply_segment_fmt(run, seg: TextSegment, base_fmt: Dict[str, Any]):
        """Apply base formatting then overlay segment-level markup."""
        if base_fmt["name"] is not None:
            run.font.name = base_fmt["name"]
        if base_fmt["size"] is not None:
            run.font.size = base_fmt["size"]
        if base_fmt["underline"] is not None:
            run.font.underline = base_fmt["underline"]

        # bold / italic: markup overrides base when True, else fall back to base
        run.font.bold = True if seg.bold else base_fmt["bold"]
        run.font.italic = True if seg.italic else base_fmt["italic"]

        # font size override from markup
        if seg.size is not None:
            run.font.size = Pt(seg.size)

        # color: markup takes priority, then base
        if seg.color is not None:
            try:
                r = int(seg.color[0:2], 16)
                g = int(seg.color[2:4], 16)
                b = int(seg.color[4:6], 16)
                run.font.color.rgb = RGBColor(r, g, b)
            except Exception:
                pass
        elif base_fmt["color_type"] == MSO_COLOR_TYPE.RGB and base_fmt["color_rgb"] is not None:
            try:
                run.font.color.rgb = base_fmt["color_rgb"]
            except Exception:
                pass
        elif base_fmt["color_scheme_xml"] is not None:
            try:
                rPr = run._r.find(qn("a:rPr"))
                if rPr is not None:
                    existing = rPr.find(qn("a:solidFill"))
                    if existing is not None:
                        rPr.remove(existing)
                    rPr.insert(0, copy.deepcopy(base_fmt["color_scheme_xml"]))
            except Exception:
                pass

    def _expand_rich_paragraphs(self, paragraph, rich_paras: list, base_fmt: Dict[str, Any]):
        """
        Replace a single paragraph with one or more rich-text paragraphs.

        The first RichParagraph reuses the existing paragraph element;
        subsequent ones are inserted after it in the txBody.
        """
        txBody = paragraph._p.getparent()
        ref_p = paragraph._p
        insert_idx = list(txBody).index(ref_p)

        for para_idx, rich_para in enumerate(rich_paras):
            if para_idx == 0:
                # Reuse the existing paragraph element
                p_elem = ref_p
                # Clear existing runs
                paragraph.clear()
            else:
                # Create a new paragraph by deep-copying the original
                p_elem = copy.deepcopy(ref_p)
                # Remove any runs carried over from the copy
                for r in p_elem.findall(qn("a:r")):
                    p_elem.remove(r)
                insert_idx += 1
                txBody.insert(insert_idx, p_elem)

            # Apply bullet formatting
            pPr = p_elem.find(qn("a:pPr"))
            if pPr is None:
                pPr = etree.SubElement(p_elem, qn("a:pPr"))
                p_elem.insert(0, pPr)

            # Remove any existing bullet elements before (re-)applying
            for tag in (qn("a:buNone"), qn("a:buChar"), qn("a:buAutoNum")):
                for el in pPr.findall(tag):
                    pPr.remove(el)

            if rich_para.bullet:
                bu = etree.SubElement(pPr, qn("a:buChar"))
                bu.set("char", "•")
            else:
                # Explicitly suppress inherited bullets so non-bullet lines
                # inside a bullet-styled text box are rendered correctly.
                etree.SubElement(pPr, qn("a:buNone"))

            # Add a run for each segment — insert before endParaRPr if present
            end_par_rpr = p_elem.find(qn("a:endParaRPr"))
            for seg in rich_para.segments:
                # SubElement で追加してから endParaRPr の前に移動
                r_elem = etree.SubElement(p_elem, qn("a:r"))
                t_elem = etree.SubElement(r_elem, qn("a:t"))
                t_elem.text = seg.text

                if end_par_rpr is not None:
                    end_par_rpr.addprevious(r_elem)

                # Build a temporary run object to leverage python-pptx font API
                from pptx.text.text import _Run  # noqa: PLC0415
                tmp_run = _Run(r_elem, paragraph)
                self._apply_segment_fmt(tmp_run, seg, base_fmt)

    # ------------------------------------------------------------------

    def _replace_images(self, slide: Slide, images: Dict[str, dict]):
        """
        Replace ``{{ img:key }}`` marker shapes with images.

        A marker is a text shape whose text contains ``{{ img:key }}``. The image is
        placed at the marker shape's position/size and the marker shape is removed.
        Markers whose key is not present in ``images`` are removed without inserting
        anything (so no marker text leaks into the output).

        Args:
            slide: The slide to process
            images: Mapping of marker key to an image spec dict:
                - ``data`` (bytes) or ``path`` (str/Path): the image source.
                  Raster formats (PNG/JPEG/...) are placed as-is. SVG sources are
                  embedded natively (``asvg:svgBlip``) with a rasterized PNG
                  fallback; this requires the optional ``resvg-py`` dependency
                  (``pip install 'pptx-templatex[svg]'``).
                - ``fit``: ``"contain"`` (default; letterbox, no cropping) or
                  ``"cover"`` (fill the frame, center-cropped)
        """
        for shape in list(slide.shapes):
            if not getattr(shape, "has_text_frame", False):
                continue
            match = _IMG_MARKER.search(shape.text_frame.text)
            if not match:
                continue
            spec = images.get(match.group(1))
            if spec is not None:
                self._place_image(slide, shape, spec)
            shape._element.getparent().remove(shape._element)

    @classmethod
    def _place_image(cls, slide: Slide, marker_shape, spec: dict):
        """Insert the image described by ``spec`` at the marker shape's bounds."""
        data = spec.get("data")
        if data is None:
            path = spec.get("path")
            if path is None:
                raise TemplateError("image spec must have 'data' (bytes) or 'path'")
            data = Path(path).read_bytes()
        fit = spec.get("fit", "contain")

        left, top = marker_shape.left, marker_shape.top
        frame_w, frame_h = marker_shape.width, marker_shape.height

        # SVGはフォールバックPNGにラスタライズして配置し、後段でSVG本体をsvgBlipとして付与する
        svg_data = data if cls._is_svg(data) else None
        if svg_data is not None:
            data = cls._rasterize_svg(svg_data, frame_w, frame_h)

        with Image.open(io.BytesIO(data)) as image:
            image_w, image_h = image.size

        if fit == "cover":
            # 枠を満たして中央クロップ（縦横比のはみ出しぶんを crop_* で切る）
            picture = slide.shapes.add_picture(io.BytesIO(data), left, top, frame_w, frame_h)
            frame_ratio = frame_w / frame_h
            image_ratio = image_w / image_h
            if image_ratio > frame_ratio:
                crop = (1 - frame_ratio / image_ratio) / 2
                picture.crop_left = picture.crop_right = crop
            elif image_ratio < frame_ratio:
                crop = (1 - image_ratio / frame_ratio) / 2
                picture.crop_top = picture.crop_bottom = crop
        else:
            # contain: 枠内に収めて中央配置（レターボックス。内容を切らない）
            scale = min(frame_w / image_w, frame_h / image_h)
            width, height = int(image_w * scale), int(image_h * scale)
            picture_left = int(left + (frame_w - width) / 2)
            picture_top = int(top + (frame_h - height) / 2)
            picture = slide.shapes.add_picture(io.BytesIO(data), picture_left, picture_top, width, height)

        if svg_data is not None:
            cls._attach_svg(slide, picture, svg_data)

    @staticmethod
    def _is_svg(data: bytes) -> bool:
        """Return True if ``data`` looks like an SVG document."""
        return b"<svg" in data[:2048]

    @staticmethod
    def _rasterize_svg(svg_data: bytes, frame_w: int, frame_h: int) -> bytes:
        """Render ``svg_data`` to PNG bytes sized for the frame (aspect ratio preserved)."""
        try:
            import resvg_py  # noqa: PLC0415
        except ImportError:
            raise TemplateError(
                "SVG images require the 'resvg-py' package. "
                "Install it with: pip install 'pptx-templatex[svg]'"
            )
        max_px = 4096
        px_w = min(max_px, max(1, round(frame_w / _EMU_PER_PX) * _SVG_FALLBACK_SCALE))
        px_h = min(max_px, max(1, round(frame_h / _EMU_PER_PX) * _SVG_FALLBACK_SCALE))
        try:
            # width/height 両指定時は resvg がアスペクト比を保って内接サイズに丸める
            return bytes(
                resvg_py.svg_to_bytes(svg_string=svg_data.decode("utf-8"), width=px_w, height=px_h)
            )
        except Exception as e:
            raise TemplateError(f"Failed to rasterize SVG: {str(e)}")

    @staticmethod
    def _attach_svg(slide: Slide, picture, svg_data: bytes):
        """
        Attach the original SVG to ``picture`` via the ``asvg:svgBlip`` extension.

        The picture keeps its PNG blip as a fallback for viewers without SVG
        support; PowerPoint 2019 / Microsoft 365 render the SVG part instead.
        """
        package = slide.part.package
        partname = package.next_partname("/ppt/media/image%d.svg")
        svg_part = Part(partname, "image/svg+xml", package, blob=svg_data)
        rId = slide.part.relate_to(svg_part, RT.IMAGE)

        blip = picture._element.blipFill.find(qn("a:blip"))
        extLst = etree.SubElement(blip, qn("a:extLst"))
        ext = etree.SubElement(extLst, qn("a:ext"))
        ext.set("uri", _SVG_EXT_URI)
        svgBlip = etree.SubElement(ext, "{%s}svgBlip" % _SVG_NS, nsmap={"asvg": _SVG_NS})
        svgBlip.set(qn("r:embed"), rId)

    def _expand_table_rows(self, slide: Slide, table_rows: Dict[str, list]):
        """
        Expand table template rows into one row per data item.

        A "template row" is a table row whose cells contain placeholders of the
        form ``{{ <key>.<field> }}`` where ``<key>`` is a key in ``table_rows``.
        The row is duplicated once per item (inheriting borders, fill and fonts
        from the template row), each copy's placeholders are resolved against
        that item, and the template row itself is removed. An empty list removes
        the template row without inserting anything.

        Args:
            slide: The slide to process
            table_rows: Mapping of row-set key to a list of row dicts
                (e.g. ``{"rows": [{"c1": "...", "c2": "..."}, ...]}``)
        """
        from copy import deepcopy

        from pptx.table import _Cell

        for shape in slide.shapes:
            if not getattr(shape, "has_table", False):
                continue
            tbl = shape.table._tbl
            for tr in list(tbl.tr_lst):
                row_text = " ".join(
                    _Cell(tc, tr).text_frame.text for tc in tr.tc_lst
                )
                matched_key = next(
                    (
                        key
                        for key in table_rows
                        if re.search(r"\{\{\s*" + re.escape(key) + r"\.", row_text)
                    ),
                    None,
                )
                if matched_key is None:
                    continue

                parent = tr.getparent()
                insert_at = parent.index(tr)
                for offset, item in enumerate(table_rows[matched_key]):
                    new_tr = deepcopy(tr)
                    parent.insert(insert_at + offset, new_tr)
                    for tc in new_tr.tc_lst:
                        self._replace_placeholders_in_text_frame(
                            _Cell(tc, new_tr).text_frame, {matched_key: item}
                        )
                parent.remove(tr)

    def _replace_placeholders_in_slide(self, slide: Slide, replacements: Dict[str, Any]):
        """
        Replace all {{ }} placeholders in a slide's text with values.

        Traverses plain shapes' text frames and table cells' text frames.

        Args:
            slide: The slide to process
            replacements: Dictionary with replacement values
        """
        for shape in slide.shapes:
            if getattr(shape, "has_table", False):
                for row in shape.table.rows:
                    for cell in row.cells:
                        self._replace_placeholders_in_text_frame(cell.text_frame, replacements)
                continue

            if not hasattr(shape, "text_frame"):
                continue

            self._replace_placeholders_in_text_frame(shape.text_frame, replacements)

    def _replace_placeholders_in_text_frame(self, text_frame, replacements: Dict[str, Any]):
        """
        Replace all {{ }} placeholders in a text frame.

        Since PowerPoint may split placeholders across multiple runs,
        we need to process the entire paragraph text at once.

        Args:
            text_frame: The text frame to process
            replacements: Dictionary with replacement values
        """
        for paragraph in text_frame.paragraphs:
            # Get the full paragraph text
            full_text = paragraph.text

            # Check if there are any placeholders in this paragraph
            if "{{" not in full_text or "}}" not in full_text:
                continue

            # Replace placeholders in the full text
            try:
                new_text = PlaceholderReplacer.replace_text(full_text, replacements)
            except Exception:
                # If replacement fails, skip this paragraph
                continue

            # If text hasn't changed, skip
            if new_text == full_text:
                continue

            # Clean control characters from the replaced text
            # Convert vertical tab (0x0B) to newline - PowerPoint uses this for soft line breaks
            new_text = new_text.replace('\x0B', '\n')
            # Remove other control characters (except \n and \r)
            new_text = re.sub(r'[\x00-\x08\x0C\x0E-\x1F]', '', new_text)

            # Find the first run with defined formatting to use as reference
            reference_run = None
            for run in paragraph.runs:
                if run.font.name is not None:
                    reference_run = run
                    break

            # If no reference found, use the first run
            if reference_run is None and len(paragraph.runs) > 0:
                reference_run = paragraph.runs[0]

            base_fmt = self._save_run_fmt(reference_run) if reference_run else {
                "name": None, "size": None, "bold": None, "italic": None,
                "underline": None, "color_type": None, "color_rgb": None,
            }

            # --- Rich text path ---
            if is_rich_text(new_text):
                rich_paras = parse_rich_text(new_text)
                self._expand_rich_paragraphs(paragraph, rich_paras, base_fmt)
                continue

            # --- Plain text path (original behaviour) ---
            # Clear all existing runs
            paragraph.clear()

            # Create a new run with the replaced text
            new_run = paragraph.add_run()
            new_run.text = new_text
            self._apply_run_fmt(new_run, base_fmt)

    def process(
        self,
        config: Union[Dict, str, Path],
        output_path: Union[str, Path]
    ):
        """
        Process the template with the given configuration and save to output file.

        Args:
            config: Configuration dict or path to JSON config file
            output_path: Path to save the output PPTX file

        Raises:
            TemplateError: If processing fails
        """
        # Load config if it's a file path
        if isinstance(config, (str, Path)):
            config_path = Path(config)
            if not config_path.exists():
                raise TemplateError(f"Config file not found: {config}")
            try:
                with open(config_path, 'r', encoding='utf-8') as f:
                    config = json.load(f)
            except Exception as e:
                raise TemplateError(f"Failed to load config: {str(e)}")

        # Validate config
        if not isinstance(config, dict) or "slides" not in config:
            raise TemplateError("Config must contain 'slides' key")

        slides_config = config["slides"]
        if not isinstance(slides_config, list):
            raise TemplateError("'slides' must be a list")

        # Create new presentation based on template to preserve theme and layouts
        # This ensures theme fonts, colors, layouts, and other settings are maintained
        output_prs = Presentation(str(self.template_path))

        # Remove all slides from the template
        while len(output_prs.slides) > 0:
            rId = output_prs.slides._sldIdLst[0].rId
            output_prs.part.drop_rel(rId)
            del output_prs.slides._sldIdLst[0]

        # Process each slide configuration
        for idx, slide_config in enumerate(slides_config):
            if not isinstance(slide_config, dict):
                raise TemplateError(f"Slide config at index {idx} must be a dict")

            if "src_page" not in slide_config:
                raise TemplateError(f"Slide config at index {idx} missing 'src_page'")

            src_page = slide_config["src_page"]
            replace_texts = slide_config.get("replace_texts", {})
            replace_table_rows = slide_config.get("replace_table_rows", {})
            replace_images = slide_config.get("replace_images", {})
            if not isinstance(replace_images, dict) or any(
                not isinstance(spec, dict) for spec in replace_images.values()
            ):
                raise TemplateError(f"'replace_images' at index {idx} must be a dict of image spec dicts")
            if not isinstance(replace_table_rows, dict) or any(
                not isinstance(rows, list) or any(not isinstance(item, dict) for item in rows)
                for rows in replace_table_rows.values()
            ):
                raise TemplateError(
                    f"'replace_table_rows' at index {idx} must be a dict of lists of dicts"
                )

            # Validate src_page
            if not isinstance(src_page, int) or src_page < 1:
                raise TemplateError(
                    f"Invalid src_page {src_page} at index {idx}: must be positive integer"
                )

            if src_page > len(self.template_prs.slides):
                raise TemplateError(
                    f"src_page {src_page} at index {idx} exceeds template slides count "
                    f"({len(self.template_prs.slides)})"
                )

            # Copy the slide from template_prs to output_prs
            # Since output_prs is also based on the same template, layouts should match
            new_slide = self._copy_slide(src_page - 1, output_prs)

            # Note: Font normalization is no longer needed because the theme is preserved
            # from the template. Font name: None will use the theme's default font.
            # self._normalize_fonts_in_slide(new_slide)

            # Expand variable-length table rows first, then replace placeholders
            if replace_table_rows:
                self._expand_table_rows(new_slide, replace_table_rows)
            if replace_texts:
                self._replace_placeholders_in_slide(new_slide, replace_texts)
            # Resolve image markers last. Runs even when replace_images is empty so
            # unassigned {{ img:key }} markers never leak into the output.
            self._replace_images(new_slide, replace_images)

        # Save output
        try:
            output_prs.save(str(output_path))
        except Exception as e:
            raise TemplateError(f"Failed to save output: {str(e)}")
