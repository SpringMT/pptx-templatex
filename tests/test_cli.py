"""Unit tests for CLI."""

import json
import sys

import pytest
from pptx import Presentation
from pptx.util import Inches

from pptx_templatex.cli import main


@pytest.fixture
def temp_dir(tmp_path):
    """Create a temporary directory for test files."""
    return tmp_path


@pytest.fixture
def sample_template(temp_dir):
    """Create a sample PowerPoint template file."""
    template_path = temp_dir / "template.pptx"
    prs = Presentation()

    slide1 = prs.slides.add_slide(prs.slide_layouts[5])
    textbox1 = slide1.shapes.add_textbox(Inches(1), Inches(1), Inches(5), Inches(1))
    textbox1.text_frame.text = "Hello {{ name }}!"

    prs.save(str(template_path))
    return template_path


@pytest.fixture
def sample_config(temp_dir):
    """Create a sample config file."""
    config_path = temp_dir / "config.json"
    config = {
        "slides": [
            {"src_page": 1, "replace_texts": {"name": "World"}}
        ]
    }
    config_path.write_text(json.dumps(config), encoding='utf-8')
    return config_path


class TestCLI:
    """Tests for CLI functionality."""

    def test_successful_execution(self, sample_template, sample_config, temp_dir, monkeypatch):
        """Test successful CLI execution."""
        output_path = temp_dir / "output.pptx"

        # Mock sys.argv
        monkeypatch.setattr(sys, "argv", [
            "pptx-templatex",
            str(sample_template),
            str(sample_config),
            str(output_path)
        ])

        result = main()
        assert result == 0
        assert output_path.exists()

    def test_template_not_found(self, sample_config, temp_dir, monkeypatch, capsys):
        """Test error when template file doesn't exist."""
        output_path = temp_dir / "output.pptx"

        monkeypatch.setattr(sys, "argv", [
            "pptx-templatex",
            str(temp_dir / "nonexistent.pptx"),
            str(sample_config),
            str(output_path)
        ])

        result = main()
        assert result == 1

        captured = capsys.readouterr()
        assert "Template file not found" in captured.err

    def test_config_not_found(self, sample_template, temp_dir, monkeypatch, capsys):
        """Test error when config file doesn't exist."""
        output_path = temp_dir / "output.pptx"

        monkeypatch.setattr(sys, "argv", [
            "pptx-templatex",
            str(sample_template),
            str(temp_dir / "nonexistent.json"),
            str(output_path)
        ])

        result = main()
        assert result == 1

        captured = capsys.readouterr()
        assert "Config file not found" in captured.err

    def test_invalid_config_file(self, sample_template, temp_dir, monkeypatch, capsys):
        """Test error when config file is invalid."""
        config_path = temp_dir / "invalid.json"
        config_path.write_text("{ invalid json }", encoding='utf-8')
        output_path = temp_dir / "output.pptx"

        monkeypatch.setattr(sys, "argv", [
            "pptx-templatex",
            str(sample_template),
            str(config_path),
            str(output_path)
        ])

        result = main()
        assert result == 1

        captured = capsys.readouterr()
        assert "Error:" in captured.err

    def test_output_created_with_correct_content(self, sample_template, sample_config, temp_dir, monkeypatch):
        """Test that output file contains correct content."""
        output_path = temp_dir / "output.pptx"

        monkeypatch.setattr(sys, "argv", [
            "pptx-templatex",
            str(sample_template),
            str(sample_config),
            str(output_path)
        ])

        result = main()
        assert result == 0

        # Verify content
        output_prs = Presentation(str(output_path))
        assert len(output_prs.slides) == 1

        text_content = []
        for shape in output_prs.slides[0].shapes:
            if hasattr(shape, "text"):
                text_content.append(shape.text)

        text_str = " ".join(text_content)
        assert "Hello World!" in text_str
        assert "{{ name }}" not in text_str

    def test_help_message(self, monkeypatch, capsys):
        """Test help message display."""
        monkeypatch.setattr(sys, "argv", ["pptx-templatex", "--help"])

        with pytest.raises(SystemExit) as exc_info:
            main()

        assert exc_info.value.code == 0

        captured = capsys.readouterr()
        assert "PowerPoint template engine" in captured.out
        assert "Examples:" in captured.out
        assert "Placeholder syntax:" in captured.out
