"""Unit tests for TemplateEngine."""

import json

import pytest
from pptx import Presentation
from pptx.util import Inches, Pt

from pptx_templatex import TemplateEngine
from pptx_templatex.exceptions import TemplateError


@pytest.fixture
def temp_dir(tmp_path):
    """Create a temporary directory for test files."""
    return tmp_path


@pytest.fixture
def sample_template(temp_dir):
    """Create a sample PowerPoint template file."""
    template_path = temp_dir / "template.pptx"
    prs = Presentation()

    # Slide 1: Simple text with placeholders
    slide1 = prs.slides.add_slide(prs.slide_layouts[5])  # Blank layout
    textbox1 = slide1.shapes.add_textbox(Inches(1), Inches(1), Inches(5), Inches(1))
    textbox1.text_frame.text = "Hello {{ name }}!"

    # Slide 2: Multiple placeholders
    slide2 = prs.slides.add_slide(prs.slide_layouts[5])
    textbox2 = slide2.shapes.add_textbox(Inches(1), Inches(1), Inches(5), Inches(2))
    tf = textbox2.text_frame
    tf.text = "User: {{ user.name }}"
    p = tf.add_paragraph()
    p.text = "Age: {{ user.age }}"

    # Slide 3: Array access
    slide3 = prs.slides.add_slide(prs.slide_layouts[5])
    textbox3 = slide3.shapes.add_textbox(Inches(1), Inches(1), Inches(5), Inches(1))
    textbox3.text_frame.text = "First item: {{ items[0].name }}"

    prs.save(str(template_path))
    return template_path


class TestTemplateEngineInit:
    """Tests for TemplateEngine initialization."""

    def test_init_with_valid_template(self, sample_template):
        """Test initialization with valid template file."""
        engine = TemplateEngine(sample_template)
        assert engine.template_path == sample_template
        assert engine.template_prs is not None
        assert len(engine.template_prs.slides) == 3

    def test_init_with_nonexistent_file(self, temp_dir):
        """Test initialization with non-existent file."""
        with pytest.raises(TemplateError, match="Template file not found"):
            TemplateEngine(temp_dir / "nonexistent.pptx")

    def test_init_with_invalid_file(self, temp_dir):
        """Test initialization with invalid PPTX file."""
        invalid_path = temp_dir / "invalid.pptx"
        invalid_path.write_text("This is not a valid PPTX file")
        with pytest.raises(TemplateError, match="Failed to load template"):
            TemplateEngine(invalid_path)

    def test_init_with_string_path(self, sample_template):
        """Test initialization with string path."""
        engine = TemplateEngine(str(sample_template))
        assert engine.template_prs is not None


class TestTemplateEngineProcess:
    """Tests for TemplateEngine.process method."""

    def test_process_single_slide_simple_replacement(self, sample_template, temp_dir):
        """Test processing single slide with simple placeholder replacement."""
        engine = TemplateEngine(sample_template)
        output_path = temp_dir / "output.pptx"

        config = {
            "slides": [
                {"src_page": 1, "replace_texts": {"name": "World"}}
            ]
        }

        engine.process(config, output_path)

        # Verify output file was created
        assert output_path.exists()

        # Verify content
        output_prs = Presentation(str(output_path))
        assert len(output_prs.slides) == 1

        # Check that placeholder was replaced
        text_content = []
        for shape in output_prs.slides[0].shapes:
            if hasattr(shape, "text"):
                text_content.append(shape.text)

        assert any("Hello World!" in text for text in text_content)
        assert not any("{{ name }}" in text for text in text_content)

    def test_process_multiple_slides(self, sample_template, temp_dir):
        """Test processing multiple slides."""
        engine = TemplateEngine(sample_template)
        output_path = temp_dir / "output.pptx"

        config = {
            "slides": [
                {"src_page": 1, "replace_texts": {"name": "Alice"}},
                {"src_page": 2, "replace_texts": {"user": {"name": "Bob", "age": 30}}},
                {"src_page": 1, "replace_texts": {"name": "Charlie"}},
            ]
        }

        engine.process(config, output_path)

        output_prs = Presentation(str(output_path))
        assert len(output_prs.slides) == 3

    def test_process_nested_replacement(self, sample_template, temp_dir):
        """Test processing with nested key replacement."""
        engine = TemplateEngine(sample_template)
        output_path = temp_dir / "output.pptx"

        config = {
            "slides": [
                {"src_page": 2, "replace_texts": {"user": {"name": "Alice", "age": 25}}}
            ]
        }

        engine.process(config, output_path)

        output_prs = Presentation(str(output_path))
        text_content = []
        for shape in output_prs.slides[0].shapes:
            if hasattr(shape, "text"):
                text_content.append(shape.text)

        text_str = " ".join(text_content)
        assert "Alice" in text_str
        assert "25" in text_str

    def test_process_array_replacement(self, sample_template, temp_dir):
        """Test processing with array access replacement."""
        engine = TemplateEngine(sample_template)
        output_path = temp_dir / "output.pptx"

        config = {
            "slides": [
                {"src_page": 3, "replace_texts": {"items": [{"name": "Apple"}, {"name": "Banana"}]}}
            ]
        }

        engine.process(config, output_path)

        output_prs = Presentation(str(output_path))
        text_content = []
        for shape in output_prs.slides[0].shapes:
            if hasattr(shape, "text"):
                text_content.append(shape.text)

        assert any("Apple" in text for text in text_content)

    def test_process_without_replacements(self, sample_template, temp_dir):
        """Test processing slide without any replacements."""
        engine = TemplateEngine(sample_template)
        output_path = temp_dir / "output.pptx"

        config = {
            "slides": [
                {"src_page": 1}
            ]
        }

        engine.process(config, output_path)

        output_prs = Presentation(str(output_path))
        assert len(output_prs.slides) == 1

    def test_process_with_config_file(self, sample_template, temp_dir):
        """Test processing with JSON config file."""
        engine = TemplateEngine(sample_template)
        output_path = temp_dir / "output.pptx"
        config_path = temp_dir / "config.json"

        config = {
            "slides": [
                {"src_page": 1, "replace_texts": {"name": "World"}}
            ]
        }

        config_path.write_text(json.dumps(config), encoding='utf-8')
        engine.process(config_path, output_path)

        assert output_path.exists()

    def test_process_invalid_config_missing_slides(self, sample_template, temp_dir):
        """Test error when config missing 'slides' key."""
        engine = TemplateEngine(sample_template)
        output_path = temp_dir / "output.pptx"

        config = {"invalid": "config"}

        with pytest.raises(TemplateError, match="Config must contain 'slides' key"):
            engine.process(config, output_path)

    def test_process_invalid_config_slides_not_list(self, sample_template, temp_dir):
        """Test error when 'slides' is not a list."""
        engine = TemplateEngine(sample_template)
        output_path = temp_dir / "output.pptx"

        config = {"slides": "not a list"}

        with pytest.raises(TemplateError, match="'slides' must be a list"):
            engine.process(config, output_path)

    def test_process_invalid_slide_config_not_dict(self, sample_template, temp_dir):
        """Test error when slide config is not a dict."""
        engine = TemplateEngine(sample_template)
        output_path = temp_dir / "output.pptx"

        config = {"slides": ["not a dict"]}

        with pytest.raises(TemplateError, match="must be a dict"):
            engine.process(config, output_path)

    def test_process_missing_src_page(self, sample_template, temp_dir):
        """Test error when src_page is missing."""
        engine = TemplateEngine(sample_template)
        output_path = temp_dir / "output.pptx"

        config = {"slides": [{"replace_texts": {"name": "World"}}]}

        with pytest.raises(TemplateError, match="missing 'src_page'"):
            engine.process(config, output_path)

    def test_process_invalid_src_page_type(self, sample_template, temp_dir):
        """Test error when src_page is not an integer."""
        engine = TemplateEngine(sample_template)
        output_path = temp_dir / "output.pptx"

        config = {"slides": [{"src_page": "invalid"}]}

        with pytest.raises(TemplateError, match="Invalid src_page"):
            engine.process(config, output_path)

    def test_process_invalid_src_page_zero(self, sample_template, temp_dir):
        """Test error when src_page is zero."""
        engine = TemplateEngine(sample_template)
        output_path = temp_dir / "output.pptx"

        config = {"slides": [{"src_page": 0}]}

        with pytest.raises(TemplateError, match="Invalid src_page"):
            engine.process(config, output_path)

    def test_process_invalid_src_page_negative(self, sample_template, temp_dir):
        """Test error when src_page is negative."""
        engine = TemplateEngine(sample_template)
        output_path = temp_dir / "output.pptx"

        config = {"slides": [{"src_page": -1}]}

        with pytest.raises(TemplateError, match="Invalid src_page"):
            engine.process(config, output_path)

    def test_process_src_page_exceeds_template(self, sample_template, temp_dir):
        """Test error when src_page exceeds template slide count."""
        engine = TemplateEngine(sample_template)
        output_path = temp_dir / "output.pptx"

        config = {"slides": [{"src_page": 10}]}

        with pytest.raises(TemplateError, match="exceeds template slides count"):
            engine.process(config, output_path)

    def test_process_nonexistent_config_file(self, sample_template, temp_dir):
        """Test error when config file doesn't exist."""
        engine = TemplateEngine(sample_template)
        output_path = temp_dir / "output.pptx"
        config_path = temp_dir / "nonexistent.json"

        with pytest.raises(TemplateError, match="Config file not found"):
            engine.process(config_path, output_path)

    def test_process_invalid_config_file(self, sample_template, temp_dir):
        """Test error when config file is invalid JSON."""
        engine = TemplateEngine(sample_template)
        output_path = temp_dir / "output.pptx"
        config_path = temp_dir / "config.json"

        config_path.write_text("{ invalid json }", encoding='utf-8')

        with pytest.raises(TemplateError, match="Failed to load config"):
            engine.process(config_path, output_path)

    def test_process_same_slide_multiple_times(self, sample_template, temp_dir):
        """Test copying the same slide multiple times with different replacements."""
        engine = TemplateEngine(sample_template)
        output_path = temp_dir / "output.pptx"

        config = {
            "slides": [
                {"src_page": 1, "replace_texts": {"name": "First"}},
                {"src_page": 1, "replace_texts": {"name": "Second"}},
                {"src_page": 1, "replace_texts": {"name": "Third"}},
            ]
        }

        engine.process(config, output_path)

        output_prs = Presentation(str(output_path))
        assert len(output_prs.slides) == 3

        # Verify each slide has different content
        texts = []
        for slide in output_prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    texts.append(shape.text)

        text_str = " ".join(texts)
        assert "First" in text_str
        assert "Second" in text_str
        assert "Third" in text_str


class TestCopySlide:
    """Tests for _copy_slide method."""

    def test_copy_preserves_text_formatting(self, temp_dir):
        """Test that text formatting is preserved when copying."""
        template_path = temp_dir / "template.pptx"
        prs = Presentation()

        slide = prs.slides.add_slide(prs.slide_layouts[5])
        textbox = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(5), Inches(1))
        tf = textbox.text_frame
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = "Bold Text"
        run.font.bold = True
        run.font.size = Pt(24)

        prs.save(str(template_path))

        # Copy the slide
        engine = TemplateEngine(template_path)
        output_path = temp_dir / "output.pptx"
        config = {"slides": [{"src_page": 1}]}
        engine.process(config, output_path)

        # Verify formatting
        output_prs = Presentation(str(output_path))
        for shape in output_prs.slides[0].shapes:
            if hasattr(shape, "text_frame"):
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        if run.text == "Bold Text":
                            assert run.font.bold is True
                            assert run.font.size == Pt(24)

    def test_copy_multiple_text_boxes(self, temp_dir):
        """Test copying slide with multiple text boxes."""
        template_path = temp_dir / "template.pptx"
        prs = Presentation()

        slide = prs.slides.add_slide(prs.slide_layouts[5])
        textbox1 = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(2), Inches(1))
        textbox1.text_frame.text = "Box 1"
        textbox2 = slide.shapes.add_textbox(Inches(4), Inches(1), Inches(2), Inches(1))
        textbox2.text_frame.text = "Box 2"

        prs.save(str(template_path))

        engine = TemplateEngine(template_path)
        output_path = temp_dir / "output.pptx"
        config = {"slides": [{"src_page": 1}]}
        engine.process(config, output_path)

        output_prs = Presentation(str(output_path))
        texts = [shape.text for shape in output_prs.slides[0].shapes if hasattr(shape, "text")]
        text_str = " ".join(texts)
        assert "Box 1" in text_str
        assert "Box 2" in text_str

    def test_font_normalization_with_none_font_name(self, temp_dir):
        """Test that runs with None font names are normalized to inherit from neighboring runs."""
        template_path = temp_dir / "template.pptx"
        prs = Presentation()

        slide = prs.slides.add_slide(prs.slide_layouts[5])
        textbox = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(5), Inches(1))
        tf = textbox.text_frame
        p = tf.paragraphs[0]

        # Create first run with explicit font
        run1 = p.add_run()
        run1.text = "Text with "
        run1.font.name = "Arial"
        run1.font.size = Pt(12)

        # Create second run with None font name (simulates PowerPoint's behavior with special chars)
        run2 = p.add_run()
        run2.text = "{{ placeholder }}"
        # Don't set font name - it will be None

        # Create third run with explicit font
        run3 = p.add_run()
        run3.text = " text"
        run3.font.name = "Arial"
        run3.font.size = Pt(12)

        prs.save(str(template_path))

        # Process the template
        engine = TemplateEngine(template_path)
        output_path = temp_dir / "output.pptx"
        config = {
            "slides": [
                {"src_page": 1, "replace_texts": {"placeholder": "REPLACED"}}
            ]
        }
        engine.process(config, output_path)

        # Verify all runs have defined font names
        output_prs = Presentation(str(output_path))
        for shape in output_prs.slides[0].shapes:
            if hasattr(shape, "text_frame"):
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        if run.text:
                            # All runs should have a font name (not None)
                            assert run.font.name is not None
                            # Should inherit Arial from the first run
                            assert run.font.name == "Arial"

    def test_font_normalization_all_none_uses_default(self, temp_dir):
        """Test that when all runs have None font, theme font is preserved."""
        template_path = temp_dir / "template.pptx"
        prs = Presentation()

        slide = prs.slides.add_slide(prs.slide_layouts[5])
        textbox = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(5), Inches(1))
        tf = textbox.text_frame
        p = tf.paragraphs[0]

        # Create run without setting font name
        run = p.add_run()
        run.text = "{{ text }}"
        # Don't set font.name - it will be None

        prs.save(str(template_path))

        # Process the template
        engine = TemplateEngine(template_path)
        output_path = temp_dir / "output.pptx"
        config = {
            "slides": [
                {"src_page": 1, "replace_texts": {"text": "Sample"}}
            ]
        }
        engine.process(config, output_path)

        # Verify font is preserved (may be None, which means use theme font)
        output_prs = Presentation(str(output_path))
        for shape in output_prs.slides[0].shapes:
            if hasattr(shape, "text_frame"):
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        if run.text:
                            # Font name can be None, which means use theme's default font
                            # This is valid and will be rendered correctly by PowerPoint
                            assert run.text == "Sample"


@pytest.fixture
def table_template(temp_dir):
    """Create a template with a native table (header row + template row)."""
    template_path = temp_dir / "table_template.pptx"
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    frame = slide.shapes.add_table(2, 3, Inches(1), Inches(1), Inches(8), Inches(2))
    table = frame.table
    for col, head in enumerate(("{{ head1 }}", "{{ head2 }}", "Result")):
        table.cell(0, col).text = head
    for col in range(3):
        table.cell(1, col).text = f"{{{{ rows.c{col + 1} }}}}"
    prs.save(str(template_path))
    return template_path


def _get_table(pptx_path):
    prs = Presentation(str(pptx_path))
    for shape in prs.slides[0].shapes:
        if getattr(shape, "has_table", False):
            return shape.table
    raise AssertionError("no table found")


class TestTemplateEngineTableRows:
    """Tests for table cell replacement and variable-row expansion."""

    def test_replace_texts_reaches_table_cells(self, table_template, temp_dir):
        """Header placeholders inside table cells are replaced via replace_texts."""
        engine = TemplateEngine(table_template)
        output = temp_dir / "out.pptx"
        engine.process(
            {"slides": [{"src_page": 1, "replace_texts": {"head1": "Event", "head2": "Action"}}]},
            output,
        )
        table = _get_table(output)
        assert table.cell(0, 0).text == "Event"
        assert table.cell(0, 1).text == "Action"

    def test_expand_table_rows(self, table_template, temp_dir):
        """The template row is duplicated once per data item."""
        engine = TemplateEngine(table_template)
        output = temp_dir / "out.pptx"
        engine.process(
            {
                "slides": [
                    {
                        "src_page": 1,
                        "replace_table_rows": {
                            "rows": [
                                {"c1": "a1", "c2": "a2", "c3": "a3"},
                                {"c1": "b1", "c2": "b2", "c3": "b3"},
                                {"c1": "c1v", "c2": "c2v", "c3": "c3v"},
                            ]
                        },
                    }
                ]
            },
            output,
        )
        table = _get_table(output)
        rows = list(table.rows)
        assert len(rows) == 4  # header + 3 data rows
        assert table.cell(1, 0).text == "a1"
        assert table.cell(2, 1).text == "b2"
        assert table.cell(3, 2).text == "c3v"

    def test_expand_table_rows_empty_list_removes_template_row(self, table_template, temp_dir):
        """An empty list removes the template row entirely."""
        engine = TemplateEngine(table_template)
        output = temp_dir / "out.pptx"
        engine.process(
            {"slides": [{"src_page": 1, "replace_table_rows": {"rows": []}}]},
            output,
        )
        table = _get_table(output)
        assert len(list(table.rows)) == 1  # header only

    def test_expand_table_rows_missing_key_leaves_row(self, table_template, temp_dir):
        """A template row whose key is not provided is left untouched."""
        engine = TemplateEngine(table_template)
        output = temp_dir / "out.pptx"
        engine.process(
            {"slides": [{"src_page": 1, "replace_table_rows": {"other": [{"x": "1"}]}}]},
            output,
        )
        table = _get_table(output)
        assert len(list(table.rows)) == 2
        assert table.cell(1, 0).text == "{{ rows.c1 }}"

    def test_expand_table_rows_rich_text_in_cell(self, table_template, temp_dir):
        """Rich text markup works inside expanded cells."""
        engine = TemplateEngine(table_template)
        output = temp_dir / "out.pptx"
        engine.process(
            {
                "slides": [
                    {
                        "src_page": 1,
                        "replace_table_rows": {
                            "rows": [{"c1": "**bold** plain", "c2": "x", "c3": "y"}]
                        },
                    }
                ]
            },
            output,
        )
        table = _get_table(output)
        cell = table.cell(1, 0)
        runs = cell.text_frame.paragraphs[0].runs
        assert cell.text == "bold plain"
        assert runs[0].font.bold is True

    def test_expand_table_rows_inherits_template_row_format(self, table_template, temp_dir):
        """Expanded rows keep the template row's cell fill (format inheritance)."""
        from pptx.dml.color import RGBColor as _RGB

        # 事前にテンプレート行へ塗りを付けた版を作る
        prs = Presentation(str(table_template))
        table = None
        for shape in prs.slides[0].shapes:
            if getattr(shape, "has_table", False):
                table = shape.table
        cell = table.cell(1, 0)
        cell.fill.solid()
        cell.fill.fore_color.rgb = _RGB(0x11, 0x22, 0x33)
        colored = table_template.parent / "colored.pptx"
        prs.save(str(colored))

        engine = TemplateEngine(colored)
        output = table_template.parent / "out.pptx"
        engine.process(
            {
                "slides": [
                    {
                        "src_page": 1,
                        "replace_table_rows": {
                            "rows": [{"c1": "a", "c2": "b", "c3": "c"}, {"c1": "d", "c2": "e", "c3": "f"}]
                        },
                    }
                ]
            },
            output,
        )
        out_table = _get_table(output)
        for row_idx in (1, 2):
            assert out_table.cell(row_idx, 0).fill.fore_color.rgb == _RGB(0x11, 0x22, 0x33)

    def test_invalid_replace_table_rows_raises(self, table_template, temp_dir):
        """Non list-of-dicts values are rejected."""
        engine = TemplateEngine(table_template)
        with pytest.raises(TemplateError, match="replace_table_rows"):
            engine.process(
                {"slides": [{"src_page": 1, "replace_table_rows": {"rows": "not-a-list"}}]},
                temp_dir / "out.pptx",
            )


@pytest.fixture
def image_template(temp_dir):
    """Create a template with {{ img:key }} marker shapes."""
    template_path = temp_dir / "image_template.pptx"
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    marker1 = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(4), Inches(3))
    marker1.text_frame.text = "{{ img:photo1 }}"
    marker2 = slide.shapes.add_textbox(Inches(6), Inches(1), Inches(3), Inches(3))
    marker2.text_frame.text = "{{ img:photo2 }}"
    prs.save(str(template_path))
    return template_path


def _png_bytes(width=200, height=100):
    import io as _io

    from PIL import Image as PILImage

    buffer = _io.BytesIO()
    PILImage.new("RGB", (width, height), "navy").save(buffer, format="PNG")
    return buffer.getvalue()


def _pictures_and_texts(pptx_path):
    prs = Presentation(str(pptx_path))
    pictures, texts = [], []
    for shape in prs.slides[0].shapes:
        if shape.shape_type == 13:  # PICTURE
            pictures.append(shape)
        elif getattr(shape, "has_text_frame", False) and shape.text_frame.text.strip():
            texts.append(shape.text_frame.text)
    return pictures, texts


class TestTemplateEngineImages:
    """Tests for {{ img:key }} image marker embedding."""

    def test_contain_fits_inside_marker_bounds(self, image_template, temp_dir):
        engine = TemplateEngine(image_template)
        output = temp_dir / "out.pptx"
        engine.process(
            {"slides": [{"src_page": 1, "replace_images": {"photo1": {"data": _png_bytes(), "fit": "contain"}}}]},
            output,
        )
        pictures, texts = _pictures_and_texts(output)
        assert len(pictures) == 1
        picture = pictures[0]
        # 200x100 (2:1) を 4x3 インチ枠に contain → 幅いっぱい・高さは半分（枠内・切らない）
        assert picture.width == Inches(4)
        assert picture.height == Inches(2)
        assert picture.top == Inches(1) + Inches(0.5)  # 縦中央
        assert not any("img:" in t for t in texts)  # マーカーは残らない

    def test_cover_fills_and_crops(self, image_template, temp_dir):
        engine = TemplateEngine(image_template)
        output = temp_dir / "out.pptx"
        engine.process(
            {"slides": [{"src_page": 1, "replace_images": {"photo2": {"data": _png_bytes(), "fit": "cover"}}}]},
            output,
        )
        pictures, _ = _pictures_and_texts(output)
        assert len(pictures) == 1
        picture = pictures[0]
        assert picture.width == Inches(3) and picture.height == Inches(3)  # 枠を満たす
        assert picture.crop_left > 0 and picture.crop_left == picture.crop_right  # 中央クロップ

    def test_fit_defaults_to_contain(self, image_template, temp_dir):
        engine = TemplateEngine(image_template)
        output = temp_dir / "out.pptx"
        engine.process(
            {"slides": [{"src_page": 1, "replace_images": {"photo1": {"data": _png_bytes()}}}]},
            output,
        )
        pictures, _ = _pictures_and_texts(output)
        assert pictures[0].crop_left == 0  # contain は切らない

    def test_unassigned_markers_are_removed(self, image_template, temp_dir):
        engine = TemplateEngine(image_template)
        output = temp_dir / "out.pptx"
        engine.process({"slides": [{"src_page": 1}]}, output)  # replace_images なし
        pictures, texts = _pictures_and_texts(output)
        assert pictures == []
        assert not any("img:" in t for t in texts)

    def test_image_from_path(self, image_template, temp_dir):
        image_path = temp_dir / "photo.png"
        image_path.write_bytes(_png_bytes())
        engine = TemplateEngine(image_template)
        output = temp_dir / "out.pptx"
        engine.process(
            {"slides": [{"src_page": 1, "replace_images": {"photo1": {"path": str(image_path)}}}]},
            output,
        )
        pictures, _ = _pictures_and_texts(output)
        assert len(pictures) == 1

    def test_same_key_resolved_per_slide(self, image_template, temp_dir):
        """同一テンプレページの複数回コピーでも slide 単位の指定で別画像になる。"""
        engine = TemplateEngine(image_template)
        output = temp_dir / "out.pptx"
        engine.process(
            {
                "slides": [
                    {"src_page": 1, "replace_images": {"photo1": {"data": _png_bytes(200, 100)}}},
                    {"src_page": 1, "replace_images": {"photo1": {"data": _png_bytes(100, 200)}}},
                ]
            },
            output,
        )
        prs = Presentation(str(output))
        sizes = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if shape.shape_type == 13:
                    sizes.append((shape.width, shape.height))
        assert len(sizes) == 2
        assert sizes[0] != sizes[1]  # 縦横比の違いが配置に反映されている

    def test_invalid_spec_raises(self, image_template, temp_dir):
        engine = TemplateEngine(image_template)
        with pytest.raises(TemplateError, match="replace_images"):
            engine.process(
                {"slides": [{"src_page": 1, "replace_images": {"photo1": "not-a-dict"}}]},
                temp_dir / "out.pptx",
            )
        with pytest.raises(TemplateError, match="data.*path|'data'"):
            engine.process(
                {"slides": [{"src_page": 1, "replace_images": {"photo1": {"fit": "cover"}}}]},
                temp_dir / "out.pptx",
            )
