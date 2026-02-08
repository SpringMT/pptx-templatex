#!/usr/bin/env python3
"""Analyze text formatting in slides."""

import sys
from pptx import Presentation

if len(sys.argv) < 2:
    print("Usage: python3 analyze_text_format.py <file.pptx>")
    sys.exit(1)

pptx_path = sys.argv[1]
print(f"Analyzing: {pptx_path}\n")

prs = Presentation(pptx_path)

for slide_idx, slide in enumerate(prs.slides, 1):
    print(f"=== Slide {slide_idx} ===")

    for shape_idx, shape in enumerate(slide.shapes, 1):
        if not hasattr(shape, "text_frame"):
            continue

        if not shape.text_frame.text.strip():
            continue

        print(f"\nShape {shape_idx}: {shape.text_frame.text[:50]}")

        for para_idx, paragraph in enumerate(shape.text_frame.paragraphs):
            if not paragraph.text.strip():
                continue

            print(f"  Paragraph {para_idx+1}:")

            for run_idx, run in enumerate(paragraph.runs):
                if not run.text:
                    continue

                print(f"    Run {run_idx+1}: '{run.text}'")
                print(f"      Font name: {run.font.name}")
                print(f"      Font size: {run.font.size}")
                print(f"      Bold: {run.font.bold}")
                print(f"      Italic: {run.font.italic}")
                print(f"      Color type: {run.font.color.type if hasattr(run.font.color, 'type') else 'N/A'}")

                if run.font.color.type == 1:  # RGB
                    try:
                        rgb = run.font.color.rgb
                        print(f"      Color RGB: {rgb}")
                    except:
                        print(f"      Color RGB: (error getting RGB)")
                elif run.font.color.type == 2:  # SCHEME
                    print(f"      Color: SCHEME")

    print()
